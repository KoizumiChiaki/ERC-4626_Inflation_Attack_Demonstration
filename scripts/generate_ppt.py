# -*- coding: utf-8 -*-
"""Generate ERC-4626 inflation attack lab presentation (.pptx)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "ERC4626通胀攻击实验汇报.pptx"

ACCENT = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_run(run, size=18, bold=False, color=DARK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Microsoft YaHei"


def add_title_bar(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    set_run(p.runs[0], 28, True, WHITE)
    tf.margin_left = Inches(0.35)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.35), Inches(1.15), Inches(9.3), Inches(0.4))
        stf = sub.text_frame
        stf.text = subtitle
        set_run(stf.paragraphs[0].runs[0], 14, False, MUTED)


def add_bullets(slide, items, top=1.55, size=20, line_space=1.25):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(top), Inches(9.1), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(10)
        p.line_spacing = line_space
        if p.runs:
            set_run(p.runs[0], size)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT
    bg.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(1.2))
    tf = t.text_frame
    tf.text = "经典智能合约漏洞复现实验"
    set_run(tf.paragraphs[0].runs[0], 36, True, WHITE)

    s = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(8.8), Inches(0.8))
    stf = s.text_frame
    stf.text = "ERC-4626 金库通胀攻击（首存者 / 捐赠攻击）"
    set_run(stf.paragraphs[0].runs[0], 24, False, WHITE)

    meta = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(8.8), Inches(0.8))
    mtf = meta.text_frame
    mtf.text = "课程实验汇报  |  Hardhat 本地靶场复现  |  erc4626-inflation-lab"
    set_run(mtf.paragraphs[0].runs[0], 14, False, RGBColor(0xDB, 0xEA, 0xFE))


def content_slide(prs, title, bullets, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title, subtitle)
    add_bullets(slide, bullets)


def two_col_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)

    lt = slide.shapes.add_textbox(Inches(0.45), Inches(1.45), Inches(4.4), Inches(0.35))
    ltf = lt.text_frame
    ltf.text = left_title
    set_run(ltf.paragraphs[0].runs[0], 16, True, ACCENT)

    lb = slide.shapes.add_textbox(Inches(0.45), Inches(1.85), Inches(4.4), Inches(3.8))
    lbf = lb.text_frame
    for i, item in enumerate(left_items):
        p = lbf.paragraphs[0] if i == 0 else lbf.add_paragraph()
        p.text = f"• {item}"
        set_run(p.runs[0], 15)

    rt = slide.shapes.add_textbox(Inches(5.15), Inches(1.45), Inches(4.4), Inches(0.35))
    rtf = rt.text_frame
    rtf.text = right_title
    set_run(rtf.paragraphs[0].runs[0], 16, True, ACCENT)

    rb = slide.shapes.add_textbox(Inches(5.15), Inches(1.85), Inches(4.4), Inches(3.8))
    rbf = rb.text_frame
    for i, item in enumerate(right_items):
        p = rbf.paragraphs[0] if i == 0 else rbf.add_paragraph()
        p.text = f"• {item}"
        set_run(p.runs[0], 15)


def code_slide(prs, title, code_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    box = slide.shapes.add_shape(1, Inches(0.45), Inches(1.55), Inches(9.1), Inches(3.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
    box.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if p.runs:
            r = p.runs[0]
            r.font.name = "Consolas"
            r.font.size = Pt(13)
            r.font.color.rgb = DARK


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    title_slide(prs)

    content_slide(
        prs,
        "实验目标",
        [
            "选取 1 个以太坊历史上的经典漏洞类型并完成本地复现",
            "使用 Solidity 编写含漏洞的简易合约靶场",
            "编写攻击脚本，演示资金被盗窃的完整过程",
            "理解漏洞根因，并提出可行的防护思路",
        ],
        subtitle="对应课程要求 · 要求.txt",
    )

    content_slide(
        prs,
        "什么是 ERC-4626？",
        [
            "EIP-4626：代币化金库（Tokenized Vault）标准，2021 草案，2022 年起广泛采用",
            "统一 deposit / mint / withdraw / redeem 与份额换算接口",
            "份额（shares）代表对金库底层资产（assets）的比例所有权",
            "Yearn、Morpho 等大量 DeFi 协议采用，成为「收益乐高」基础模块",
        ],
    )

    content_slide(
        prs,
        "漏洞概述：通胀攻击",
        [
            "别名：首存者攻击、捐赠攻击、份额价格操纵",
            "目标：流动性极低、刚部署的空金库或近乎空的 ERC-4626 金库",
            "核心：操纵「每份份额对应的资产单价」，使后续存款向下取整为 0 份",
            "不依赖外部预言机或闪电贷，属于内部分额会计漏洞",
        ],
    )

    content_slide(
        prs,
        "攻击四步流程",
        [
            "① 首存：攻击者存入 1 wei，获得 1 share（唯一股东）",
            "② 捐赠：直接向金库 transfer 大量代币，不 mint 份额，抬高 totalAssets",
            "③ 受害者存款：调用 deposit，整数除法舍入 → 获得 0 shares，资产却已进入金库",
            "④ 攻击者赎回：redeem 全部份额，取走池中资产（含受害者资金）",
        ],
    )

    code_slide(
        prs,
        "关键公式（VulnerableVault）",
        [
            "totalAssets() = asset.balanceOf(address(this))   // 含捐赠",
            "",
            "convertToShares(assets):",
            "    return (assets * totalSupply) / totalAssets   // 向下取整",
            "",
            "deposit: 若 shares == 0，仍 transferFrom 收资产，但不 mint",
        ],
    )

    content_slide(
        prs,
        "漏洞根因",
        [
            "totalAssets 使用 balanceOf，捐赠可操纵份额单价",
            "Solidity 整数除法向下取整，小额存款可能得 0 shares",
            "deposit 缺少 minSharesOut 校验，0 份仍收走资产",
            "缺少虚拟份额 / 死份额，空金库无最低 exchange rate 锚定",
            "攻击者可 frontrun 受害者首笔存款",
        ],
    )

    two_col_slide(
        prs,
        "与预言机攻击的对比",
        "预言机 + 闪电贷（2020 年代常见）",
        [
            "操纵外部价格源",
            "依赖 AMM / Oracle 喂价",
            "常配合闪电贷放大资金",
        ],
        "ERC-4626 通胀攻击",
        [
            "操纵金库内部 exchange rate",
            "totalAssets / totalSupply 会计",
            "首存 + 捐赠即可，成本与受害存款相关",
        ],
    )

    content_slide(
        prs,
        "典型链上案例（节选）",
        [
            "2022-03：Solmate #178 公开讨论；OpenZeppelin 2023 v4.9 引入虚拟份额缓解",
            "2025-06 ResupplyFi：cvcrvUSD 金库捐赠 + 借贷汇率算 0，约 $9.6M",
            "2025-05 Inertia：roETH 份额价格操纵，约 $152k",
            "2026-03 dTRINITY：经典首存捐赠，约 $257k 坏账",
            "说明：4626 汇率被借贷/预言机误用时，危害可远超「偷首存」",
        ],
    )

    content_slide(
        prs,
        "靶场架构",
        [
            "MockUSDC.sol — 实验用 ERC20 测试资产（mUSDC）",
            "VulnerableVault.sol — 故意含漏洞的简化 ERC-4626 金库",
            "InflationAttacker.sol — 封装 setup（首存+捐赠）与 steal（赎回）",
            "scripts/attack.js — 控制台分步演示攻击",
            "test/InflationAttack.test.js — 自动化断言受害者 0 份、攻击者获利",
        ],
        subtitle="技术栈：Hardhat 2.x + Solidity 0.8.20 + OpenZeppelin",
    )

    two_col_slide(
        prs,
        "实验参数与角色",
        "攻击参数（与脚本一致）",
        [
            "首存 SEED = 1 wei",
            "捐赠 DONATION = 10,000 mUSDC",
            "受害者存入 9,999 mUSDC",
        ],
        "链上角色",
        [
            "deployer：部署合约",
            "attacker：执行攻击",
            "victim：正常用户存款",
        ],
    )

    content_slide(
        prs,
        "运行与验证",
        [
            "npm install",
            "npm run compile",
            "npm test          → Chai 断言：受害者份额 = 0，攻击者余额增加",
            "npm run attack    → 打印各步金库份额、余额与结论",
            "环境：Node.js 18+，本地 Hardhat 链，仅用于教学",
        ],
    )

    content_slide(
        prs,
        "预期实验结果",
        [
            "setup 后：金库 totalSupply = 1，totalAssets ≈ 10,000 mUSDC",
            "受害者 deposit 9,999 mUSDC 后：balanceOf(victim) = 0",
            "攻击者 steal() 后：攻击者 USDC 余额显著增加",
            "金库剩余资产接近 0，受害者无法 redeem",
            "结论：资产进库 ≠ 获得份额，信任边界在份额与资产是否同步",
        ],
    )

    content_slide(
        prs,
        "防护建议",
        [
            "虚拟份额 / 虚拟资产偏移（OpenZeppelin ERC4626 v4.9+ _decimalsOffset）",
            "部署时铸造死份额并锁死（如 mint 到 0xdead）",
            "deposit 要求 shares >= minSharesOut，否则 revert",
            "集成方用 Router 校验 previewDeposit 与实得份额",
            "借贷/预言机勿裸用 convertToAssets，需 TWAP、下限检查等",
        ],
    )

    content_slide(
        prs,
        "小结",
        [
            "本实验在本地 Hardhat 链完整复现 ERC-4626 通胀攻击资金路径",
            "漏洞本质是空金库 + balanceOf 计账 + 整数舍入 + 无滑点保护",
            "该问题 2022 年起已被广泛讨论，但 2025–2026 仍有主网损失案例",
            "开发金库或集成 4626 时，必须假设首存窗口为敌对环境",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "谢谢聆听")
    add_bullets(
        slide,
        [
            "项目仓库：erc4626-inflation-lab",
            "延伸阅读：ERC4626通胀攻击介绍.md",
            "Q & A",
        ],
        top=2.0,
        size=24,
    )

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
